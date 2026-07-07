"""Generate downloadable teaching artifacts: Word, PowerPoint, Markdown, HTML."""

from __future__ import annotations

import html
import re
from pathlib import Path


def safe_filename(value: str, fallback: str = "export") -> str:
    cleaned = re.sub(r"[^0-9A-Za-z\u4e00-\u9fff._-]+", "_", (value or fallback).strip())
    return cleaned[:80] or fallback


def write_markdown(path: Path, title: str, body: str) -> None:
    content = f"# {title}\n\n{body.strip()}\n"
    path.write_text(content, encoding="utf-8")


def write_docx(path: Path, title: str, body: str, sections: list[dict] | None = None) -> None:
    from docx import Document
    from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
    from docx.shared import Pt

    doc = Document()
    heading = doc.add_heading(title, level=0)
    heading.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
    doc.add_paragraph(f"生成时间：自动生成 · AI错题拆博士")
    doc.add_paragraph("")

    if sections:
        for section in sections:
            section_title = section.get("title") or section.get("heading") or "章节"
            doc.add_heading(str(section_title), level=1)
            content = section.get("content") or section.get("body") or ""
            for block in str(content).split("\n"):
                block = block.strip()
                if not block:
                    continue
                if block.startswith("- "):
                    doc.add_paragraph(block[2:], style="List Bullet")
                elif re.match(r"^\d+\.\s", block):
                    doc.add_paragraph(re.sub(r"^\d+\.\s", "", block), style="List Number")
                else:
                    doc.add_paragraph(block)
            for question in section.get("questions") or []:
                if isinstance(question, dict):
                    stem = question.get("stem") or question.get("title") or ""
                    answer = question.get("answer") or ""
                    analysis = question.get("analysis") or ""
                    if stem:
                        doc.add_paragraph(str(stem))
                    if answer:
                        p = doc.add_paragraph()
                        p.add_run("【答案】").bold = True
                        p.add_run(str(answer))
                    if analysis:
                        p = doc.add_paragraph()
                        p.add_run("【解析】").bold = True
                        p.add_run(str(analysis))
                else:
                    doc.add_paragraph(str(question))
    else:
        for block in str(body).split("\n"):
            block = block.strip()
            if not block:
                continue
            if block.startswith("## "):
                doc.add_heading(block[3:], level=1)
            elif block.startswith("### "):
                doc.add_heading(block[4:], level=2)
            elif block.startswith("- "):
                doc.add_paragraph(block[2:], style="List Bullet")
            elif re.match(r"^\d+\.\s", block):
                doc.add_paragraph(re.sub(r"^\d+\.\s", "", block), style="List Number")
            else:
                doc.add_paragraph(block)

    for paragraph in doc.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Microsoft YaHei"
            run.font.size = Pt(11)
    doc.save(path)


def write_pptx(path: Path, title: str, slides: list[dict]) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = title
    try:
        slide.placeholders[1].text = "AI错题拆博士 · 讲评课件"
    except (KeyError, IndexError):
        pass

    bullet_layout = prs.slide_layouts[1]
    for item in slides:
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = str(item.get("title") or "讲评页")
        body = slide.placeholders[1].text_frame
        body.clear()
        bullets = item.get("bullets") or item.get("points") or []
        if not bullets and item.get("content"):
            bullets = [line.strip() for line in str(item.get("content")).split("\n") if line.strip()]
        for idx, bullet in enumerate(bullets[:8]):
            paragraph = body.paragraphs[0] if idx == 0 else body.add_paragraph()
            paragraph.text = str(bullet)
            paragraph.level = 0
            paragraph.font.size = Pt(20)
        notes = item.get("speaker_notes") or item.get("notes") or ""
        if notes:
            slide.notes_slide.notes_text_frame.text = str(notes)
    prs.save(path)


def write_knowledge_map_html(path: Path, title: str, nodes: list[dict], edges: list[dict] | None = None) -> None:
    edges = edges or []
    node_html = []
    for idx, node in enumerate(nodes):
        label = html.escape(str(node.get("label") or node.get("name") or f"节点{idx + 1}"))
        detail = html.escape(str(node.get("detail") or node.get("description") or ""))
        node_html.append(
            f"""
            <div class="node">
              <strong>{label}</strong>
              <p>{detail}</p>
            </div>
            """
        )
    edge_html = []
    for edge in edges:
        src = html.escape(str(edge.get("from") or edge.get("source") or ""))
        dst = html.escape(str(edge.get("to") or edge.get("target") or ""))
        label = html.escape(str(edge.get("label") or ""))
        if src and dst:
            edge_html.append(f"<li><span>{src}</span> → <span>{dst}</span>{f' · {label}' if label else ''}</li>")

    page = f"""<!doctype html>
<html lang="zh-CN">
<head>
  <meta charset="UTF-8" />
  <title>{html.escape(title)}</title>
  <style>
    body {{ font-family: "Microsoft YaHei", sans-serif; background:#f7f3ea; color:#1f1a24; padding:24px; }}
    h1 {{ text-align:center; }}
    .grid {{ display:grid; grid-template-columns:repeat(auto-fit,minmax(220px,1fr)); gap:16px; margin-top:24px; }}
    .node {{ background:#fff; border:1px solid #e4d6c2; border-radius:10px; padding:14px; box-shadow:0 8px 20px rgba(0,0,0,.06); }}
    .node strong {{ display:block; margin-bottom:8px; color:#67131d; }}
    .edges {{ margin-top:24px; background:#fff; border-radius:10px; padding:16px 24px; }}
    .edges li {{ margin:8px 0; }}
  </style>
</head>
<body>
  <h1>{html.escape(title)}</h1>
  <div class="grid">{''.join(node_html)}</div>
  {'<div class="edges"><h2>关系链</h2><ul>' + ''.join(edge_html) + '</ul></div>' if edge_html else ''}
</body>
</html>
"""
    path.write_text(page, encoding="utf-8")


def artifact_record(filename: str, url: str, kind: str, label: str) -> dict:
    return {
        "filename": filename,
        "url": url,
        "kind": kind,
        "label": label,
    }
